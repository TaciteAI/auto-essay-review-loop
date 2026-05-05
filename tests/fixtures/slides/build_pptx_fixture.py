#!/usr/bin/env python3
"""
build_pptx_fixture.py - generate a tiny .pptx fixture for edit_pptx.py tests.

Built on demand (not committed as a binary blob) so the fixture stays
diffable. Uses python-pptx if available; exits 0 with a noop message
if not (the test harness then skips the round-trip test).

The fixture deck has 3 slides:
  Slide 1: Title "Q4 strategy" + body "We should focus on three things." (no notes)
  Slide 2: Title "Problem" + body containing a long sentence to be cut +
           speaker notes already populated
  Slide 3: Title "Ask" + body "Hire two engineers."

Usage:
    py tests/fixtures/slides/build_pptx_fixture.py <output_path>
"""
from __future__ import annotations

import sys
from pathlib import Path


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print("usage: build_pptx_fixture.py <output_path>", file=sys.stderr)
        return 2
    out_path = Path(argv[1])

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("python-pptx not installed; skipping fixture build", file=sys.stderr)
        return 0  # not an error; the test harness handles skip

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[5]  # "Title Only"

    # Slide 1
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Q4 strategy"
    tx1 = s1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf1 = tx1.text_frame
    tf1.text = "We should focus on three things."

    # Slide 2 — has both a wordy bullet line and pre-existing notes
    s2 = prs.slides.add_slide(title_layout)
    s2.shapes.title.text = "Problem"
    tx2 = s2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(3))
    tf2 = tx2.text_frame
    tf2.text = "Customers waste an average of $432K per year on duplicated work and abandoned projects."
    p = tf2.add_paragraph()
    p.text = "This is the wordy descriptor we plan to cut."
    s2.notes_slide.notes_text_frame.text = "Original notes for slide 2."

    # Slide 3
    s3 = prs.slides.add_slide(title_layout)
    s3.shapes.title.text = "Ask"
    tx3 = s3.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf3 = tx3.text_frame
    tf3.text = "Hire two engineers."

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"built fixture: {out_path}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
