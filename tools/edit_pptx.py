#!/usr/bin/env python3
"""
edit_pptx.py - apply a structured fix list to a .pptx in place.

The auto-slides-review-loop skill emits prose fix lists by default; for
.pptx input it also emits a STRUCTURED fix list (JSON) that this writer
consumes. The writer mutates a copy of the deck via python-pptx so the
user's PowerPoint-authored visual design (theme, layout, fonts, accent
colors) is preserved. Only text content and speaker notes are touched.

Optional dependency: python-pptx >= 0.6.21. If the package is missing,
the writer fails fast with a JSON error so the loop can fall back to
the prose fix list. The companion verifier (verify_slides.py) stays
stdlib-only.

Usage (from a skill via the run.sh wrapper):
    bash tools/run.sh edit_pptx.py \\
        --input  review-stage/draft.pptx \\
        --output review-stage/draft_round2.pptx \\
        --fixes  review-stage/fix_list_round_1.json

Direct invocation:
    py tools/edit_pptx.py --input X.pptx --output Y.pptx --fixes Z.json

Fix list schema (top-level):
    {
      "schema_version": 1,
      "source": "review-stage/draft.pptx",
      "fixes": [ <fix-object>, ... ]
    }

Each fix object:
    {
      "id":           "<short stable id, e.g. 'density-f1'>",
      "slide_index":  <1-based slide number>,
      "action":       "set_notes" | "append_notes" |
                      "replace_text" | "delete_text" | "replace_title",
      "target":       "<text to find> (required for replace_text / delete_text)",
      "replacement":  "<new text> (required for set_notes / append_notes /
                       replace_text / replace_title)",
      "rationale":    "<short reason; logged but not applied>"
    }

Actions:
    set_notes      Replace the slide's speaker notes wholesale.
    append_notes   Append a paragraph to the slide's existing notes
                   (creating the notes slide if needed).
    replace_text   Find `target` text inside any text frame on the slide
                   and replace it with `replacement`. Match is exact and
                   case-sensitive after whitespace normalization.
    delete_text    Like replace_text but with replacement="". Empty
                   paragraphs left behind are removed.
    replace_title  Replace the title placeholder text on the slide.

Out-of-scope actions (will be reported as `skipped` with reason
"unsupported in v0.2"): add_slide, delete_slide, move_slide, restyle,
add_image. The skill surfaces those to the user as manual fixes.

Output (stdout, JSON):
    {
      "tool": "edit_pptx",
      "schema_version": 1,
      "timestamp": "...",
      "input_file":  "review-stage/draft.pptx",
      "output_file": "review-stage/draft_round2.pptx",
      "fixes_total": 7,
      "applied":  [{"id":"...", "slide_index":3, "action":"...", "result":"applied"}, ...],
      "skipped":  [{"id":"...", "action":"add_slide", "reason":"unsupported in v0.2"}, ...],
      "warnings": [{"id":"...", "reason":"target text not found on slide 3"}, ...],
      "passed":   true,
      "summary":  "5 of 7 fixes applied (1 skipped, 1 warning)"
    }

Exit codes:
    0  success — at least one fix was applied OR all fixes were no-ops/skipped cleanly
    1  partial — at least one fix produced a warning (target not found, etc.)
    2  hard error — bad input, missing dep, unparseable .pptx
"""
from __future__ import annotations

import json
import sys
from datetime import datetime, timezone
from pathlib import Path

TOOL_NAME = "edit_pptx"
SCHEMA_VERSION = 1

# Hard-skip actions: out of scope for v0.2 but recognized so the skill can
# surface them as manual TODOs without crashing.
UNSUPPORTED_ACTIONS = {
    "add_slide", "delete_slide", "move_slide", "reorder_slides",
    "restyle", "add_image", "add_shape", "change_layout",
}
SUPPORTED_ACTIONS = {
    "set_notes", "append_notes",
    "replace_text", "delete_text", "replace_title",
}


def now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def emit_error(input_file: str, output_file: str, summary: str, exc_name: str) -> None:
    err = {
        "tool": TOOL_NAME,
        "schema_version": SCHEMA_VERSION,
        "timestamp": now_iso(),
        "input_file": input_file,
        "output_file": output_file,
        "fixes_total": 0,
        "applied": [],
        "skipped": [],
        "warnings": [],
        "passed": False,
        "summary": summary,
        "error": exc_name,
    }
    sys.stdout.write(json.dumps(err, ensure_ascii=False, indent=2))
    sys.stdout.write("\n")


# ---------------------------------------------------------------------------
# Argument parsing
# ---------------------------------------------------------------------------

def parse_args(argv: list[str]) -> dict:
    out = {"input": "", "output": "", "fixes": "", "dry_run": False}
    it = iter(argv[1:])
    for arg in it:
        if arg == "--dry-run":
            out["dry_run"] = True
            continue
        for prefix, key in (
            ("--input=", "input"),
            ("--output=", "output"),
            ("--fixes=", "fixes"),
        ):
            if arg.startswith(prefix):
                out[key] = arg[len(prefix):]
                break
        else:
            # Allow space-separated form: --input X.pptx
            if arg in ("--input", "--output", "--fixes"):
                try:
                    out[arg.lstrip("-")] = next(it)
                except StopIteration as e:
                    raise SystemExit(f"missing value for {arg}") from e
            else:
                raise SystemExit(f"unknown argument: {arg}")
    if not out["input"] or not out["output"] or not out["fixes"]:
        raise SystemExit(
            "usage: edit_pptx.py --input X.pptx --output Y.pptx --fixes Z.json [--dry-run]"
        )
    return out


# ---------------------------------------------------------------------------
# python-pptx helpers
# ---------------------------------------------------------------------------

def _normalize_ws(s: str) -> str:
    """Collapse runs of whitespace to a single space, strip ends. Used so a
    target like '  hello  world  ' matches 'hello world' inside a text frame
    without forcing the persona to reproduce smart-quote / soft-hyphen exotica."""
    return " ".join((s or "").split())


def _iter_text_frames(slide):
    """Yield every text_frame on the slide that holds non-title content.

    We yield title frames too — `replace_text` legitimately targets title
    text — and let the action helpers decide whether to skip them based on
    placeholder type.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        yield shape


def _shape_is_title(shape) -> bool:
    """Best-effort title detection. python-pptx exposes is_placeholder + ph type."""
    try:
        if not shape.is_placeholder:
            return False
        ph = shape.placeholder_format
        # ph.type is an enum; 13 = TITLE, 14 = CTR_TITLE, 15 = SUBTITLE.
        # Use the underlying int to avoid importing PP_PLACEHOLDER (less brittle
        # across python-pptx versions).
        return ph.idx == 0 or (ph.type is not None and int(ph.type) in (13, 14))
    except (AttributeError, ValueError):
        return False


def _replace_in_text_frame(tf, target: str, replacement: str) -> bool:
    """Replace `target` -> `replacement` inside a text frame, preserving as much
    formatting as possible.

    Strategy:
      1. If a single run contains the target, replace within that run (best:
         keeps font/color/size).
      2. Else if a single paragraph contains the target across runs, write the
         replacement into the first run and clear the rest of the paragraph's
         runs (preserves the paragraph's first-run formatting).
      3. Else if the whole frame's text contains the target across paragraphs,
         skip — we don't rewrite multi-paragraph spans (too easy to clobber).

    Returns True if a replacement was made.
    """
    target_norm = _normalize_ws(target)
    if not target_norm:
        return False

    # 1. Per-run match.
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            if target_norm in _normalize_ws(run.text or ""):
                # Preserve formatting; only rewrite the .text.
                # If the run's normalized text equals the target, replace whole.
                if _normalize_ws(run.text) == target_norm:
                    run.text = replacement
                else:
                    # Substring match: replace the substring, keeping surrounding text.
                    run.text = _replace_substring_preserving_outer(run.text, target, replacement)
                return True

    # 2. Per-paragraph match.
    for paragraph in tf.paragraphs:
        full = "".join(r.text or "" for r in paragraph.runs)
        if target_norm in _normalize_ws(full):
            # Build the new paragraph text.
            new_text = _replace_substring_preserving_outer(full, target, replacement)
            # Write into the first run; clear the rest. This preserves the
            # paragraph-level alignment + the first-run font.
            runs = list(paragraph.runs)
            if not runs:
                paragraph.text = new_text
            else:
                runs[0].text = new_text
                for run in runs[1:]:
                    run.text = ""
            return True

    return False


def _replace_substring_preserving_outer(full: str, target: str, replacement: str) -> str:
    """Replace the first occurrence of target (whitespace-tolerant) inside full,
    keeping characters outside the match unchanged.

    We do whitespace-normalized matching but reconstruct the original outside
    spans byte-for-byte. This keeps trailing punctuation and surrounding
    whitespace intact.
    """
    target_norm = _normalize_ws(target)
    if not target_norm:
        return full

    # Try plain string match first (covers 95% of cases with no surprises).
    if target in full:
        return full.replace(target, replacement, 1)

    # Whitespace-tolerant match: collapse a copy, find positions, map back.
    collapsed_chars: list[str] = []
    index_map: list[int] = []  # index_map[i] = original index of collapsed_chars[i]
    prev_was_space = False
    for orig_i, ch in enumerate(full):
        if ch.isspace():
            if prev_was_space:
                continue
            collapsed_chars.append(" ")
            index_map.append(orig_i)
            prev_was_space = True
        else:
            collapsed_chars.append(ch)
            index_map.append(orig_i)
            prev_was_space = False
    collapsed = "".join(collapsed_chars).strip()
    # `collapsed` may have a leading space chopped; track that offset.
    leading_strip = 0
    if collapsed_chars and collapsed_chars[0] == " ":
        leading_strip = 1

    pos = collapsed.find(target_norm)
    if pos == -1:
        return full  # no match, unchanged
    start_collapsed = pos + leading_strip
    end_collapsed = start_collapsed + len(target_norm)
    # Map collapsed indices back to original-string indices.
    if start_collapsed >= len(index_map):
        return full
    start_orig = index_map[start_collapsed]
    if end_collapsed - 1 >= len(index_map):
        return full
    end_orig = index_map[end_collapsed - 1] + 1
    return full[:start_orig] + replacement + full[end_orig:]


def _delete_in_text_frame(tf, target: str) -> bool:
    """Delete `target` from a text frame, removing emptied paragraphs."""
    if not _replace_in_text_frame(tf, target, ""):
        return False
    # Clean up: drop paragraphs whose text is now blank.
    # python-pptx requires deleting via the underlying XML.
    paras_xml = tf._txBody.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
    )
    for p in paras_xml:
        text = "".join(
            (t.text or "")
            for t in p.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
        )
        if not text.strip():
            # Don't delete the only remaining paragraph (python-pptx requires >=1).
            if len(paras_xml) > 1:
                tf._txBody.remove(p)
                paras_xml = [pp for pp in paras_xml if pp is not p]
    return True


# ---------------------------------------------------------------------------
# Per-action handlers
# ---------------------------------------------------------------------------

def apply_set_notes(slide, fix: dict) -> tuple[str, str | None]:
    replacement = fix.get("replacement", "")
    notes_slide = slide.notes_slide  # creates one if missing
    notes_slide.notes_text_frame.text = replacement
    return "applied", None


def apply_append_notes(slide, fix: dict) -> tuple[str, str | None]:
    replacement = fix.get("replacement", "")
    if not replacement:
        return "applied", "empty replacement; nothing appended"
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    existing = (tf.text or "").rstrip()
    if existing:
        tf.text = existing + "\n\n" + replacement
    else:
        tf.text = replacement
    return "applied", None


def apply_replace_text(slide, fix: dict) -> tuple[str, str | None]:
    target = fix.get("target", "")
    replacement = fix.get("replacement", "")
    if not target:
        return "warning", "missing target"
    for shape in _iter_text_frames(slide):
        if _replace_in_text_frame(shape.text_frame, target, replacement):
            return "applied", None
    return "warning", f"target text not found on slide {fix.get('slide_index')}"


def apply_delete_text(slide, fix: dict) -> tuple[str, str | None]:
    target = fix.get("target", "")
    if not target:
        return "warning", "missing target"
    for shape in _iter_text_frames(slide):
        if _delete_in_text_frame(shape.text_frame, target):
            return "applied", None
    return "warning", f"target text not found on slide {fix.get('slide_index')}"


def apply_replace_title(slide, fix: dict) -> tuple[str, str | None]:
    replacement = fix.get("replacement", "")
    for shape in _iter_text_frames(slide):
        if _shape_is_title(shape):
            shape.text_frame.text = replacement
            return "applied", None
    return "warning", f"no title placeholder found on slide {fix.get('slide_index')}"


ACTION_HANDLERS = {
    "set_notes": apply_set_notes,
    "append_notes": apply_append_notes,
    "replace_text": apply_replace_text,
    "delete_text": apply_delete_text,
    "replace_title": apply_replace_title,
}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def load_fix_list(path: Path) -> list[dict]:
    raw = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(raw, dict):
        raise ValueError("fix list root must be a JSON object")
    fixes = raw.get("fixes")
    if not isinstance(fixes, list):
        raise ValueError("fix list must contain a 'fixes' array")
    return fixes


def main(argv: list[str]) -> int:
    try:
        opts = parse_args(argv)
    except SystemExit as exc:
        emit_error("", "", str(exc).strip() or "argument error", "ArgumentError")
        return 2

    input_path = Path(opts["input"])
    output_path = Path(opts["output"])
    fixes_path = Path(opts["fixes"])

    if not input_path.exists():
        emit_error(str(input_path), str(output_path), f"input file not found: {input_path}", "FileNotFoundError")
        return 2
    if not fixes_path.exists():
        emit_error(str(input_path), str(output_path), f"fix list not found: {fixes_path}", "FileNotFoundError")
        return 2

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        emit_error(
            str(input_path),
            str(output_path),
            "python-pptx is not installed; install with `pip install python-pptx>=0.6.21` "
            "or run the loop with a markdown sidecar instead of .pptx",
            "ImportError",
        )
        return 2

    try:
        fixes = load_fix_list(fixes_path)
    except (json.JSONDecodeError, ValueError, OSError) as exc:
        emit_error(str(input_path), str(output_path), f"could not read fix list: {exc}", type(exc).__name__)
        return 2

    try:
        prs = Presentation(str(input_path))
    except Exception as exc:  # python-pptx raises various exceptions; keep generic.
        emit_error(str(input_path), str(output_path), f"could not open .pptx: {exc}", type(exc).__name__)
        return 2

    n_slides = len(prs.slides)
    applied: list[dict] = []
    skipped: list[dict] = []
    warnings: list[dict] = []

    for fix in fixes:
        fid = fix.get("id", "")
        action = fix.get("action", "")
        slide_idx = fix.get("slide_index")

        if action in UNSUPPORTED_ACTIONS:
            skipped.append({
                "id": fid,
                "action": action,
                "slide_index": slide_idx,
                "reason": "unsupported in v0.2 (surface as manual fix)",
            })
            continue
        if action not in SUPPORTED_ACTIONS:
            skipped.append({
                "id": fid,
                "action": action or "<missing>",
                "slide_index": slide_idx,
                "reason": f"unknown action '{action}'",
            })
            continue

        if not isinstance(slide_idx, int) or slide_idx < 1 or slide_idx > n_slides:
            warnings.append({
                "id": fid,
                "action": action,
                "slide_index": slide_idx,
                "reason": f"slide_index out of range (deck has {n_slides} slides)",
            })
            continue

        slide = prs.slides[slide_idx - 1]
        handler = ACTION_HANDLERS[action]
        try:
            result, note = handler(slide, fix)
        except Exception as exc:  # never crash on a single fix; record + continue
            warnings.append({
                "id": fid,
                "action": action,
                "slide_index": slide_idx,
                "reason": f"{type(exc).__name__}: {exc}",
            })
            continue

        if result == "applied":
            applied.append({
                "id": fid,
                "slide_index": slide_idx,
                "action": action,
                "result": "applied",
                **({"note": note} if note else {}),
            })
        else:  # warning
            warnings.append({
                "id": fid,
                "action": action,
                "slide_index": slide_idx,
                "reason": note or "unknown failure",
            })

    if not opts["dry_run"]:
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
        except (OSError, PermissionError) as exc:
            emit_error(str(input_path), str(output_path), f"could not save output: {exc}", type(exc).__name__)
            return 2

    n_total = len(fixes)
    n_applied = len(applied)
    n_skipped = len(skipped)
    n_warn = len(warnings)
    # passed = no hard errors and at least one fix was either applied or
    # cleanly skipped. An empty fix list also counts as passed (no-op).
    passed = (n_warn == 0)
    parts = [f"{n_applied} of {n_total} fixes applied"]
    if n_skipped:
        parts.append(f"{n_skipped} skipped")
    if n_warn:
        parts.append(f"{n_warn} warning(s)")
    summary = " (".join([parts[0], ", ".join(parts[1:]) + ")"]) if len(parts) > 1 else parts[0]

    result = {
        "tool": TOOL_NAME,
        "schema_version": SCHEMA_VERSION,
        "timestamp": now_iso(),
        "input_file": str(input_path),
        "output_file": str(output_path),
        "fixes_total": n_total,
        "applied": applied,
        "skipped": skipped,
        "warnings": warnings,
        "passed": passed,
        "summary": summary,
        "dry_run": opts["dry_run"],
    }
    sys.stdout.write(json.dumps(result, ensure_ascii=False, indent=2))
    sys.stdout.write("\n")
    return 0 if passed else 1


if __name__ == "__main__":
    sys.exit(main(sys.argv))
